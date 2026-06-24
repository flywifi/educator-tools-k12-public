#!/usr/bin/env python3
"""Office authoring engine (capability-gated) — produce REAL .pptx / .docx / .xlsx from a structured spec.

This is the "above baseline" deliverable: not a description of a deck, an actual editable file. Uses
python-pptx / python-docx / openpyxl WHEN installed (tools/requirements-office.txt); when a library is
absent it writes the spec JSON sidecar and returns an honest capability gap — it never emits a fake/empty
binary. Optional LibreOffice (`soffice`) conversion turns the result into PDF/PNG for visual QA
(render_convert capability). The caller (e.g. presentation-builder) adds the governed metadata block +
human_review_required; the spec itself is decision support.

CLI:
  python3 shared/office/office_authoring.py --type pptx --spec deck.json --out out/deck.pptx [--pdf]
"""
from __future__ import annotations

import argparse
import importlib.util
import json
import shutil
import subprocess
import sys
from pathlib import Path
from typing import Optional


def _have(mod: str) -> bool:
    try:
        return importlib.util.find_spec(mod) is not None
    except Exception:
        return False


def _gap(out: Path, spec: dict, capability: str) -> dict:
    sidecar = out.with_suffix(out.suffix + ".spec.json")
    sidecar.write_text(json.dumps(spec, indent=2), encoding="utf-8")
    return {"status": "capability_unavailable", "capability": capability,
            "wrote_spec": str(sidecar),
            "note": f"{capability} library not installed; wrote the spec instead of a fake file — "
                    "install tools/requirements-office.txt to emit the real document"}


def build_pptx(spec: dict, out: Path) -> dict:
    if not _have("pptx"):
        return _gap(out, spec, "office_authoring:python-pptx")
    from pptx import Presentation
    prs = Presentation()
    if spec.get("title"):
        s = prs.slides.add_slide(prs.slide_layouts[0])
        s.shapes.title.text = spec["title"]
        if spec.get("subtitle") and len(s.placeholders) > 1:
            s.placeholders[1].text = spec["subtitle"]
    for sl in spec.get("slides", []):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = sl.get("title", "")
        bullets = sl.get("bullets", [])
        if bullets and len(slide.placeholders) > 1:
            tf = slide.placeholders[1].text_frame
            for i, b in enumerate(bullets):
                (tf.paragraphs[0] if i == 0 else tf.add_paragraph()).text = str(b)
        if sl.get("notes"):
            slide.notes_slide.notes_text_frame.text = str(sl["notes"])
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    return {"status": "ok", "path": str(out), "slides": len(spec.get("slides", [])) + (1 if spec.get("title") else 0)}


def build_docx(spec: dict, out: Path) -> dict:
    if not _have("docx"):
        return _gap(out, spec, "office_authoring:python-docx")
    from docx import Document
    doc = Document()
    if spec.get("title"):
        doc.add_heading(spec["title"], 0)
    for sec in spec.get("sections", []):
        if sec.get("heading"):
            doc.add_heading(str(sec["heading"]), int(sec.get("level", 1)))
        for p in sec.get("paragraphs", []):
            doc.add_paragraph(str(p))
        for tbl in sec.get("tables", []):
            rows = tbl.get("rows", [])
            if rows:
                t = doc.add_table(rows=len(rows), cols=max(len(r) for r in rows))
                for ri, row in enumerate(rows):
                    for ci, cell in enumerate(row):
                        t.rows[ri].cells[ci].text = str(cell)
    out.parent.mkdir(parents=True, exist_ok=True)
    doc.save(str(out))
    return {"status": "ok", "path": str(out), "sections": len(spec.get("sections", []))}


def build_xlsx(spec: dict, out: Path) -> dict:
    if not _have("openpyxl"):
        return _gap(out, spec, "office_authoring:openpyxl")
    from openpyxl import Workbook
    wb = Workbook()
    first = True
    for sh in spec.get("sheets", [{"name": "Sheet1", "rows": spec.get("rows", [])}]):
        ws = wb.active if first else wb.create_sheet()
        ws.title = sh.get("name", "Sheet")
        first = False
        for row in sh.get("rows", []):
            ws.append(list(row))
    out.parent.mkdir(parents=True, exist_ok=True)
    wb.save(str(out))
    return {"status": "ok", "path": str(out), "sheets": len(spec.get("sheets", [])) or 1}


def convert(path: Path, to: str = "pdf", outdir: Optional[Path] = None) -> dict:
    """Convert/render via LibreOffice headless (render_convert capability) for visual QA."""
    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if not soffice:
        return {"status": "capability_unavailable", "capability": "render_convert",
                "note": "LibreOffice (soffice) not on PATH; cannot convert/render"}
    outdir = outdir or path.parent
    try:
        subprocess.run([soffice, "--headless", "--convert-to", to, "--outdir", str(outdir), str(path)],
                       capture_output=True, timeout=120, check=True)
        return {"status": "ok", "out": str(outdir / (path.stem + "." + to))}
    except Exception as exc:
        return {"status": "error", "detail": f"{exc.__class__.__name__}: {exc}"}


BUILDERS = {"pptx": build_pptx, "docx": build_docx, "xlsx": build_xlsx}


def main(argv) -> int:
    ap = argparse.ArgumentParser(description="Author a real Office document from a spec (capability-gated).")
    ap.add_argument("--type", choices=list(BUILDERS), required=True)
    ap.add_argument("--spec", required=True, help="path to a JSON spec")
    ap.add_argument("--out", required=True, help="output file path")
    ap.add_argument("--pdf", action="store_true", help="also render a PDF via LibreOffice (QA)")
    a = ap.parse_args(argv)
    spec = json.loads(Path(a.spec).read_text(encoding="utf-8"))
    out = Path(a.out)
    result = BUILDERS[a.type](spec, out)
    if a.pdf and result.get("status") == "ok":
        result["pdf"] = convert(out, "pdf")
    print(json.dumps(result, indent=2))
    return 0 if result.get("status") == "ok" else 0  # gap is a valid, honest outcome


if __name__ == "__main__":
    raise SystemExit(main(sys.argv[1:]))
